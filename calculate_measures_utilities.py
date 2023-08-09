import os
import sys
import tempfile
from datetime import datetime

from tqdm.contrib.concurrent import process_map
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from functools import partial
from glob import glob
from typing import List, Tuple, Dict, Optional
from skimage.morphology import binary_dilation, ball
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import nibabel as nib
from matplotlib.ticker import PercentFormatter
# sys.path.append('/home/rochman/Documents/src/Errors_Characterization')
from Errors_Characterization.utils import (write_to_excel, load_nifti_data, get_tumors_intersections,
                                           get_connected_components, dice, min_distance, pre_process_segmentation,
                                           getLargestCC, approximate_diameter, bbox2_3D)


def aggregate_pairwise_results_per_FU(ex_fn: str, output_ex_fn: str, agg: str = 'mean'):
    """
    Aggregates `calculate measures` pairwise results per Followup case.

    Parameters
    ----------
    ex_fn : str
        Path to original pairwise `calculate measures` results Excel file.
    output_ex_fn : str
        Path to Excel file where to output the aggregated results.
    agg : {'mean', 'median'}
        Aggregation method to use: 'mean' or 'median'.

    Returns
    -------
    None
    """

    writer = pd.ExcelWriter(output_ex_fn, engine='xlsxwriter')

    # read original diameters sheets
    for d in [0, 5, 10]:
        df = pd.read_excel(ex_fn, sheet_name=f'diameter_{d}').iloc[:-5]

        # aggregate per FU
        df['FU'] = [s[1] for s in df[df.columns[0]].str.split('_FU_')]
        if agg == 'mean':
            df = df.groupby('FU')[df.columns[1:-1]].mean()
        elif agg == 'median':
            df = df.groupby('FU')[df.columns[1:-1]].median()
        else:
            raise ValueError(f'The `agg` parameter must be either "mean" or "median".')
        df.reset_index(names='case_name', inplace=True)

        # save results
        write_to_excel(df, writer, df.columns, 'case_name', f'diameter_{d}',
                       {'F1 Score': ('Precision', 'Recall')}
                       )

    writer.close()


def merge_excel_results_files(ex_fns: List[str], output_ex_fn: str):
    """
    Merges different datasets' `calculate measures` results.

    Parameters
    ----------
    ex_fns : list of str
        List of paths to Excel files to merge.
    output_ex_fn : str
        Path to Excel file where to output the merged results.

    Returns
    -------
    None
    """

    assert len(ex_fns) > 0

    writer = pd.ExcelWriter(output_ex_fn, engine='xlsxwriter')

    # read original diameters sheets
    for d in [0, 5, 10]:
        res_df = pd.DataFrame()

        # merge results
        for ex_fn in ex_fns:
            res_df = pd.concat([res_df, pd.read_excel(ex_fn, sheet_name=f'diameter_{d}').iloc[:-5]]).reset_index(drop=True)

        res_df.rename(columns={res_df.columns[0]: 'case_name'}, inplace=True)

        # save results
        write_to_excel(res_df, writer, res_df.columns, 'case_name', f'diameter_{d}',
                       {'F1 Score': ('Precision', 'Recall')}
                       )

    writer.close()


def analyze_per_tumor_measures(gt_tumors_pred_tumors_roi: Tuple[str, str, str],
                               region_based_inference: bool = False) -> Tuple[List[Tuple[int, float, float, float]], List[Tuple[float, float]]]:
    """
    Analyzing each GT tumor in the given case:

    - Classifying each GT tumor as TP or FN according to the given tumors' prediction file.
    - Calculating the Dice-Coefficient score for each of them according to the given tumors' prediction file.
    - Calculating each GT tumor's approximated diameter in mm.
    - Calculating the distance of each GT tumor to organ border.

    In addition, analyzing each FP Pred tumor in the given predicted case:

    - Calculating each FP Pred tumor's approximated diameter in mm.
    - Calculating the distance of each FP Pred tumor to organ border.

    Parameters
    ----------
    gt_tumors_pred_tumors_roi : tuple of three strs
        Paths to the GT tumors mask, the tumors' prediction mask and the organ ROI mask.
    region_based_inference : bool, default False
        Whether the predicted masks were inferenced using a region-based trained model.

    Returns
    -------
    results : list of tuple of int, and 3 floats and list of tuple of 2 floats
            1) A list containing for each GT tumor:
                1) An integer indicating if the tumor is a TP (1) or an FN (0); 2) the tumor's Dice-Coefficient score;
                3) the tumor's approximated diameter in mm, and; 4) the distance of the tumor to the organ's border.
            2) A list containing for each FP Pred tumor:
                1) the tumor's approximated diameter in mm, and; 2) the distance of the tumor to the organ's border.
    """

    # extract parameters
    gt_tumors_fn, pred_tumors_fn, roi_fn = gt_tumors_pred_tumors_roi

    # load ROI mask
    roi, _ = load_nifti_data(roi_fn)

    # load the tumors masks
    gt_tumors, nifti_f = load_nifti_data(gt_tumors_fn)
    pred_tumors, _ = load_nifti_data(pred_tumors_fn)

    if region_based_inference:
        gt_tumors = (gt_tumors == 2).astype(gt_tumors.dtype)
        pred_tumors = (pred_tumors == 2).astype(pred_tumors.dtype)

    # preprocess the GT tumors masks (1)
    gt_tumors = pre_process_segmentation(gt_tumors)

    # preprocess the ROI mask
    roi = np.logical_or(roi, gt_tumors).astype(roi.dtype)
    roi = getLargestCC(roi)
    roi = pre_process_segmentation(roi, remove_small_obs=False)

    # preprocess the GT tumors masks (2)
    gt_tumors = np.logical_and(roi, gt_tumors).astype(gt_tumors.dtype)

    # crop to ROI
    xmin, xmax, ymin, ymax, zmin, zmax = bbox2_3D(roi)
    s = np.s_[max(0, xmin - 2): min(roi.shape[0], xmax + 3),
        max(0, ymin - 2): min(roi.shape[1], ymax + 3),
        max(0, zmin - 2): min(roi.shape[2], zmax + 3)]
    roi = roi[s]
    gt_tumors = gt_tumors[s]
    pred_tumors = pred_tumors[s]

    # preprocess the Pred tumors masks
    pred_tumors = np.logical_and(roi, pred_tumors).astype(pred_tumors.dtype)
    pred_tumors = pre_process_segmentation(pred_tumors)

    roi_border = np.logical_xor(roi, binary_dilation(roi, ball(1)))

    # label tumors
    gt_tumors = get_connected_components(gt_tumors, connectivity=None)
    pred_tumors = get_connected_components(pred_tumors, connectivity=None)

    # extract all the GT tumors' labels
    all_GT_tumors = np.unique(gt_tumors)
    all_GT_tumors = all_GT_tumors[all_GT_tumors != 0]

    # extract all the pred tumors' labels
    all_pred_tumors = np.unique(pred_tumors)
    all_pred_tumors = all_pred_tumors[all_pred_tumors != 0]

    # extract for each TP tumor, a list of all pred tumors intersecting with it
    intersections = get_tumors_intersections(gt_tumors, pred_tumors)

    pred_tp_tumors = list(set(t for ts in intersections.values() for t in ts))

    voxelspacing = nifti_f.header.get_zooms()

    gt_res = []

    # iterating over all GT tumors
    for current_t in all_GT_tumors:

        gt_current_t = gt_tumors == current_t

        diameter = approximate_diameter(gt_current_t.sum() * voxelspacing[0] * voxelspacing[1] * voxelspacing[2])

        # if current_t is a TP tumor:
        if current_t in intersections:
            tp = 1
            dce = dice(gt_current_t, np.isin(pred_tumors, intersections[current_t]))

        # if current_t is an FN tumor:
        else:
            tp = 0
            dce = 0

        dist = min_distance(gt_current_t, roi_border, voxelspacing=voxelspacing)

        gt_res.append((tp, dce, diameter, dist))

    pred_fp_res = []

    # iterating over all PRED tumors
    for current_t in all_pred_tumors:

        if current_t not in pred_tp_tumors:

            pred_current_t = pred_tumors == current_t

            diameter = approximate_diameter(pred_current_t.sum() * voxelspacing[0] * voxelspacing[1] * voxelspacing[2])

            dist = min_distance(pred_current_t, roi_border, voxelspacing=voxelspacing)

            pred_fp_res.append((diameter, dist))

    return gt_res, pred_fp_res


def write_per_tumor_analysis(gt_tumors_dir: str, pred_tumors_dir: str, rois_dir: str,
                             region_based_inference: bool = False):

    """
    Write per tumor analysis into an Excel file.

    Parameters
    ----------
    gt_tumors_dir : str
        Path to GT tumors directory.
    pred_tumors_dir : str
        Path to PRED tumors directory.
    rois_dir : str
        Path to ROIs directory.
    region_based_inference : bool, default False
        Whether the predicted masks were inferenced using a region-based trained model.

    Returns
    -------
    None
    """

    gt_tumors_fns = sorted(glob(f'{gt_tumors_dir}/*.nii.gz'))
    pred_tumors_fns = sorted(glob(f'{pred_tumors_dir}/*.nii.gz'))
    roi_fns = sorted(glob(f'{rois_dir}/*.nii.gz'))
    output_ex_fn = f'{pred_tumors_dir}/per_tumor_analysis.xlsx'

    # n_cases = 3 # todo delete
    n_cases = None
    gt_tumors_fns, pred_tumors_fns, roi_fns = gt_tumors_fns[:n_cases], pred_tumors_fns[:n_cases], roi_fns[:n_cases]

    case_ids = []
    for i, gt_t_fn in enumerate(gt_tumors_fns):
        assert os.path.basename(gt_t_fn) == os.path.basename(pred_tumors_fns[i])
        assert os.path.basename(gt_t_fn) == os.path.basename(roi_fns[i])
        case_ids.append(os.path.basename(gt_t_fn).replace('.nii.gz', ''))

    res = process_map(partial(analyze_per_tumor_measures, region_based_inference=region_based_inference),
                      list(zip(gt_tumors_fns, pred_tumors_fns, roi_fns)), max_workers=5)
    # res = list(map(partial(analyze_per_tumor_measures, region_based_inference=region_based_inference), list(zip(gt_tumors_fns, pred_tumors_fns, roi_fns))))
    res_gts, res_fp_preds = zip(*res)

    res_gts = [(f'{case_ids[j]}_-_{i + 1}', *r_t) for j, r in enumerate(res_gts) for i, r_t in enumerate(r)]
    gts_columns = ['Case-ID', 'Is TP?', 'Dice', 'Tumor Diameter (mm)', 'Distance from organ border (mm)']
    res_gts = pd.DataFrame(res_gts, columns=gts_columns)

    os.makedirs(os.path.dirname(output_ex_fn), exist_ok=True)
    writer = pd.ExcelWriter(output_ex_fn, engine='xlsxwriter')
    write_to_excel(res_gts, writer, gts_columns, 'Case-ID', 'GT Lesions')

    res_fp_preds = [(f'{case_ids[j]}_-_{i + 1}', *r_t) for j, r in enumerate(res_fp_preds) for i, r_t in enumerate(r)]
    fp_preds_columns = ['Case-ID', 'Tumor Diameter (mm)', 'Distance from organ border (mm)']
    res_fp_preds = pd.DataFrame(res_fp_preds, columns=fp_preds_columns)

    write_to_excel(res_fp_preds, writer, fp_preds_columns, 'Case-ID', 'FP Pred Lesions')

    writer.close()


def build_comparison_presentation_slide(data: List[List[List[List[str]]]], models: List[str], datasets: List[str],
                                        categories: List[str], measures: List[str], title: Optional[str] = None,
                                        prs: Optional[Presentation] = None) -> Presentation:

    """
    Builds models comparison presentation slide.

    Parameters
    ----------
    data : List[List[List[List[str]]]]
        The input data for the table in the slide. It is structured as a nested list, where each level represents a
        model, dataset, category, and a list of measure values, respectively.
    models : list of str
        List of model names.
    datasets : list of str
        List of dataset names.
    categories : list of str
        List of category names.
    measures : list of str
        List of measure names.
    title : str, optional
        Title to add to the slide.
    prs : Presentation object, optional
        Presentation to add the slide to. If is not given a new Presentation object is initialized.

    Returns
    -------
    prs : Presentation object
        The presentation the slide was added to. If `prs` is given, the same object will be returned.
    """

    # extract number of models
    n_models = len(data)
    assert len(models) == n_models

    # extract number of datasets
    n_datasets = len(data[0]) if n_models > 0 else 0
    assert len(datasets) == n_datasets

    # extract number of categories
    n_categories = len(data[0][0]) if n_datasets > 0 else 0
    assert len(categories) == n_categories

    # extract number of measures
    n_measures = len(data[0][0][0]) if n_categories > 0 else 0
    assert len(measures) == n_measures

    # initialize presentation
    prs = Presentation() if prs is None else prs

    # add 1 slice to the presentation
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # prepare table size and place
    x_to_place, y_to_place, cell_width, cell_height = Inches(0), Inches(0), Inches(10), Inches(1.5)
    n_rows = 2 + n_measures * n_datasets
    n_columns = 2 + n_categories * n_models

    # add slide title
    if title is not None:
        tit = slide.shapes.title
        tit.text = title
        y_to_place += tit.top + tit.height

    # add table to slide
    shape = slide.shapes.add_table(n_rows, n_columns, x_to_place, y_to_place, cell_width, cell_height)
    table = shape.table

    def write_to_table(r, c, text):
        table.cell(r, c).text = text

    # insert datasets and measures names
    table.cell(0, 0).merge(table.cell(1, 0))
    table.cell(0, 1).merge(table.cell(1, 1))
    table.cell(0, 1).text = 'Measure'
    for d, dataset in enumerate(datasets):
        table.cell(2 + n_measures * d, 0).merge(table.cell(1 + n_measures * (d + 1), 0))
        write_to_table(2 + n_measures * d, 0, dataset)
        for me, meas in enumerate(measures):
            write_to_table(2 + n_measures * d + me, 1, meas)

    # iterate over all models
    for m, model in enumerate(models):

        # insert model name
        table.cell(0, 2 + n_categories * m).merge(table.cell(0, 1 + n_categories * (m + 1)))
        write_to_table(0, 2 + n_categories * m, model)

        # insert categories names
        for c, cat in enumerate(categories):
            write_to_table(1, 2 + n_categories * m + c, cat)

        # insert measures values
        for d in range(n_datasets):
            for c in range(n_categories):
                for me in range(n_measures):
                    write_to_table(2 + n_measures * d + me, 2 + n_categories * m + c, data[m][d][c][me])

    # adjust table size and font
    font_size = 8 if n_models == 5 else 10

    for row in table.rows:
        for cell in row.cells:
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)

    return prs


def build_per_tumor_comparison_slides(ex_data: Dict[str, Dict[str, str]], prs: Optional[Presentation] = None) -> Presentation:
    """
    Builds per tumor models comparison presentation slides.

    Parameters
    ----------
    ex_data : Dict[str, Dict[str, str]]
        A dictionary containing performance data for different models and datasets.
        The keys of the `ex_data` dictionary represent model names, while the values are dictionaries.
        Each nested dictionary contains dataset names as keys and the corresponding path to an Excel file as the value.
        The Excel file at each path contains the "per tumors" performance of the model on the respective dataset.
    prs : Presentation object, optional
        Presentation to add the slide to. If is not given a new Presentation object is initialized.

    Returns
    -------
    prs : Presentation object
        The presentation the slide was added to. If `prs` is given, the same object will be returned.
    """

    def filter_diameters(_df, _min=0, _max=np.inf):
        return _df[(_min < _df['Tumor Diameter (mm)']) & (_df['Tumor Diameter (mm)'] <= _max)]

    possible_cats = {
        'all': (0, np.inf),
        '>5': (5, np.inf),
        '<5': (0, 5),
        '>10': (10, np.inf),
        '<10': (0, 10),
        '5<&<10': (5, 10),
    }

    def get_structured_data(categories):
        data = []
        for m in sorted(ex_data):
            model = []
            for d in sorted(ex_data[m]):
                ex = ex_data[m][d]

                gt_df = pd.read_excel(ex, sheet_name='GT Lesions').iloc[:-5]
                fp_pred_df = pd.read_excel(ex, sheet_name='FP Pred Lesions').iloc[:-5]

                dataset = []
                for c in categories:
                    cat_min, cat_max = possible_cats[c]
                    gt_df_to_analyse = filter_diameters(gt_df, cat_min, cat_max)
                    fp_pred_df_to_analyse = filter_diameters(fp_pred_df, cat_min, cat_max)

                    # calculating number of lesions
                    n_lesions = len(gt_df_to_analyse)

                    # calculate Dice
                    tp_dice = gt_df_to_analyse[gt_df_to_analyse['Is TP?'] == 1]['Dice']
                    mean_dice = f'{tp_dice.mean():.2f}\n+-{tp_dice.std():.2f}'

                    # calculate detection measures
                    n_tp = len(tp_dice)
                    n_fn = n_lesions - n_tp
                    n_fp = len(fp_pred_df_to_analyse)
                    precision = f'{n_tp / (n_tp + n_fp):.2f}'
                    recall = f'{n_tp / (n_tp + n_fn):.2f}'

                    dataset.append([f'{n_lesions}', mean_dice, precision, recall, f'{n_tp}', f'{n_fn}', f'{n_fp}'])

                model.append(dataset)

            data.append(model)

        return data

    categories1 = ['all', '>5', '>10']
    categories2 = ['<5', '5<&<10', '>10']

    data1 = get_structured_data(categories1)
    data2 = get_structured_data(categories2)

    models = sorted(ex_data)
    datasets = sorted(ex_data[models[0]])
    measures = ['# of lesions', 'Mean Dice', 'Precision', 'Recall', '#TP', '#FN', '#FP']

    prs = build_comparison_presentation_slide(data1, models, datasets, categories1, measures,
                                              title='Per Tumor Comparison (1)', prs=prs)
    build_comparison_presentation_slide(data2, models, datasets, categories2, measures,
                                        title='Per Tumor Comparison (2)', prs=prs)

    return prs


def build_per_scan_comparison_slides(ex_data: Dict[str, Dict[str, str]], prs: Optional[Presentation] = None) -> Presentation:
    """
    Builds per scan models comparison presentation slides.

    Parameters
    ----------
    ex_data : Dict[str, Dict[str, str]]
        A dictionary containing performance data for different models and datasets.
        The keys of the `ex_data` dictionary represent model names, while the values are dictionaries.
        Each nested dictionary contains dataset names as keys and the corresponding path to an Excel file as the value.
        The Excel file at each path contains the "per scan" performance of the model on the respective dataset.
    prs : Presentation object, optional
        Presentation to add the slide to. If is not given a new Presentation object is initialized.

    Returns
    -------
    prs : Presentation object
        The presentation the slide was added to. If `prs` is given, the same object will be returned.
    """

    possible_cats = {
        'all': 'diameter_0',
        '>5': 'diameter_5',
        '<5': 'diameter<5',
        '>10': 'diameter_10',
        '<10': 'diameter_<10',
        '5<&<10': 'diameter_in_(5,10)',
    }

    def get_structured_data(categories):
        data = []
        for m in sorted(ex_data):
            model = []
            for d in sorted(ex_data[m]):
                ex = ex_data[m][d]
                dataset = []
                for c in categories:
                    cat_df = pd.read_excel(ex, sheet_name=possible_cats[c]).iloc[:-5]

                    # calculating number of lesions
                    n_lesions = f'{int(cat_df["Num of lesion"].sum())}'

                    # calculate Dice
                    dice = f'{cat_df["Dice"].mean():.2f}\n+-{cat_df["Dice"].std():.2f}'

                    # calculate precision
                    precision = f'{cat_df["Precision"].mean():.2f}\n+-{cat_df["Precision"].std():.2f}'
                    recall = f'{cat_df["Recall"].mean():.2f}\n+-{cat_df["Recall"].std():.2f}'

                    dataset.append([n_lesions, dice, precision, recall])

                model.append(dataset)

            data.append(model)

        return data

    categories1 = ['all', '>5', '>10']
    categories2 = ['<5', '5<&<10', '>10']

    data1 = get_structured_data(categories1)
    data2 = get_structured_data(categories2)

    models = sorted(ex_data)
    datasets = sorted(ex_data[models[0]])
    measures = ['# of lesions', 'Dice', 'Precision', 'Recall']

    prs = build_comparison_presentation_slide(data1, models, datasets, categories1, measures,
                                              title='Per Scan Comparison (1)', prs=prs)
    build_comparison_presentation_slide(data2, models, datasets, categories2, measures,
                                        title='Per Scan Comparison (2)', prs=prs)

    return prs


if __name__ == '__main__':

    # ex_data_per_scan = {
    #     '1: Simultaneous-Duplicate - RBT\n1K epochs\n(14)': {
    #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/predLabelsStandaloneDuplicateTs_final/tumors_measurements_-_th_1.xlsx',
    #     },
    #     '2: Simultaneous-Duplicate - RBT\n2K epochs FS\n(14)': {
    #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/2000_epochs/predLabelsStandaloneDuplicateTs_best/tumors_measurements_-_th_1.xlsx',
    #     },
    #     '3: Simultaneous-Duplicate - RBT\n2K epochs PW\nLR = 5e-3\n(14)': {
    #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/2000_epochs_PW_lr_5en3/predLabelsStandaloneDuplicateTs_final/tumors_measurements_-_th_1.xlsx',
    #     },
    #     '4: Simultaneous-Duplicate - RBT\n2K epochs CL\nLR = 2.5e-3\n(14)': {
    #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/2000_epochs_CL_lr_25en4/predLabelsStandaloneDuplicateTs_final/tumors_measurements_-_th_1.xlsx',
    #     },
    #     '5: Simultaneous-Duplicate - RBT\n2K epochs CL\nLR = 1e-3\n(14)': {
    #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/2000_epochs_CL_lr_1en3/predLabelsStandaloneDuplicateTs_final/tumors_measurements_-_th_1.xlsx',
    #     },
    # }
    #
    # # ex_data_per_scan = {
    # #     '1: Simultaneous-Duplicate - RBT\n2K epochs FS\n(14)': {
    # #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/2000_epochs/predLabelsStandaloneDuplicateTs_best/tumors_measurements_-_th_1.xlsx',
    # #         '2: Val Set\n\n9 patients\n21 scans\n609 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/2000_epochs/predLabelsStandaloneDuplicateVl_best/tumors_measurements_-_th_1.xlsx',
    # #     },
    # #     '2: Simultaneous-Duplicate\n1K epochs FS\n(14)': {
    # #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/predLabelsStandaloneDuplicateTs_final/tumors_measurements_-_th_1.xlsx',
    # #         '2: Val Set\n\n9 patients\n21 scans\n609 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/predLabelsStandaloneDuplicateVl_final/tumors_measurements_-_th_1.xlsx',
    # #     }
    # # }
    #
    # prs = build_per_scan_comparison_slides(ex_data_per_scan)
    #
    # ex_data_per_tumor = {
    #     '1: Simultaneous-Duplicate - RBT\n1K epochs\n(14)': {
    #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/predLabelsStandaloneDuplicateTs_final/per_tumor_analysis.xlsx',
    #     },
    #     '2: Simultaneous-Duplicate - RBT\n2K epochs FS\n(14)': {
    #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/2000_epochs/predLabelsStandaloneDuplicateTs_best/per_tumor_analysis.xlsx',
    #     },
    #     '3: Simultaneous-Duplicate - RBT\n2K epochs PW\nLR = 5e-3\n(14)': {
    #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/2000_epochs_PW_lr_5en3/predLabelsStandaloneDuplicateTs_final/per_tumor_analysis.xlsx',
    #     },
    #     '4: Simultaneous-Duplicate - RBT\n2K epochs CL\nLR = 2.5e-3\n(14)': {
    #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/2000_epochs_CL_lr_25en4/predLabelsStandaloneDuplicateTs_final/per_tumor_analysis.xlsx',
    #     },
    #     '5: Simultaneous-Duplicate - RBT\n2K epochs CL\nLR = 1e-3\n(14)': {
    #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/2000_epochs_CL_lr_1en3/predLabelsStandaloneDuplicateTs_final/per_tumor_analysis.xlsx',
    #     },
    # }
    #
    # # ex_data_per_tumor = {
    # #     '1: Simultaneous-Duplicate - RBT\n2K epochs FS\n(14)': {
    # #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/2000_epochs/predLabelsStandaloneDuplicateTs_best/per_tumor_analysis.xlsx',
    # #         '2: Val Set\n\n9 patients\n21 scans\n609 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/2000_epochs/predLabelsStandaloneDuplicateVl_best/per_tumor_analysis.xlsx',
    # #     },
    # #     '2: Simultaneous-Duplicate\n1K epochs FS\n(14)': {
    # #         '1: Test Set\n\n17 patients\n42 scans\n667 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/predLabelsStandaloneDuplicateTs_final/per_tumor_analysis.xlsx',
    # #         '2: Val Set\n\n9 patients\n21 scans\n609 lesions': '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/predLabelsStandaloneDuplicateVl_final/per_tumor_analysis.xlsx',
    # #     }
    # # }
    #
    # prs = build_per_tumor_comparison_slides(ex_data_per_tumor, prs)
    # dt = datetime.now()
    # prs.save(f"/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Model_comparison/14_&_1K_vS 2k_epochs_FS|PW|CL_&_test_-_{dt.__str__().split('.')[0].replace(' ', '_-_')}.pptx")
    # exit(0)

    pass

    # models = [f'm{i}' for i in range(1, 3)]
    # datasets = [f'd{i}' for i in range(1, 4)]
    # categories = [f'c{i}' for i in range(1, 5)]
    # measures = [f'me{i}' for i in range(1, 6)]
    #
    # data = []
    # for m in models:
    #     model = []
    #     for d in datasets:
    #         dataset = []
    #         for c in categories:
    #             cat = []
    #             for me in measures:
    #                 cat.append(f'{m}-{d}-{c}-{me}')
    #             dataset.append(cat)
    #         model.append(dataset)
    #     data.append(model)
    #
    # prs = build_comparison_presentation_slide(data=data,
    #                                           models=models,
    #                                           datasets=datasets,
    #                                           categories=categories,
    #                                           measures=measures)
    #
    # build_comparison_presentation_slide(data=data,
    #                                     models=models,
    #                                     datasets=datasets,
    #                                     categories=categories,
    #                                     measures=measures,
    #                                     prs=prs)
    #
    # prs.save('bla.pptx')
    #
    # exit(0)

    pass

    # experiment_to_merge = [
    #     [
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset002_HadassahOriginalTrainSetExperimentStandalone_zscore_normalization/predLabelsPisaTr_best/tumors_measurements_-_th_1.xlsx',
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset002_HadassahOriginalTrainSetExperimentStandalone_zscore_normalization/predLabelsPisaTs_best/tumors_measurements_-_th_1.xlsx',
    #     ],
    #     [
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset005_HadassahOriginalTrainSetExperimentSimultaneousImprovedRegistration_zscore_normalization/predLabelsPisaTr_checkpoint_final/tumors_measurements_-_th_1_-_mean_agg_per_FU.xlsx',
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset005_HadassahOriginalTrainSetExperimentSimultaneousImprovedRegistration_zscore_normalization/predLabelsPisaTs_checkpoint_final/tumors_measurements_-_th_1_-_mean_agg_per_FU.xlsx',
    #     ],
    #     [
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset005_HadassahOriginalTrainSetExperimentSimultaneousImprovedRegistration_zscore_normalization/predLabelsStandaloneDuplicatePisaTr_checkpoint_final/tumors_measurements_-_th_1.xlsx',
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset005_HadassahOriginalTrainSetExperimentSimultaneousImprovedRegistration_zscore_normalization/predLabelsStandaloneDuplicatePisaTs_checkpoint_final/tumors_measurements_-_th_1.xlsx',
    #     ],
    #     [
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset006_HadassahOriginalTrainSetExperimentSimultaneousWithPriorImprovedRegistration_zscore_normalization/predLabelsPisaTr_checkpoint_final/tumors_measurements_-_th_1_-_mean_agg_per_FU.xlsx',
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset006_HadassahOriginalTrainSetExperimentSimultaneousWithPriorImprovedRegistration_zscore_normalization/predLabelsPisaTs_checkpoint_final/tumors_measurements_-_th_1_-_mean_agg_per_FU.xlsx',
    #     ],
    # ]
    # for excels_to_merge in experiment_to_merge:
    #     output_ex_fn = excels_to_merge[0].replace('PisaTr', 'Pisa').replace('PisaTs', 'Pisa')
    #     assert output_ex_fn != excels_to_merge[0]
    #     os.makedirs(os.path.dirname(output_ex_fn), exist_ok=True)
    #     merge_excel_results_files(excels_to_merge, output_ex_fn=output_ex_fn)

    pass

    # for gt_tumors_dir, pred_tumors_dir, rois_dir, region_based_inference in [
    #     (
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/labelsStandaloneDuplicateTs',
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/predLabelsStandaloneDuplicateTs_final',
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/roiMasksStandaloneDuplicateTs',
    #         True,
    #     ),
    #     (
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/labelsStandaloneDuplicateVl',
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/predLabelsStandaloneDuplicateVl_final',
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/roiMasksStandaloneDuplicateVl',
    #         True,
    #     ),
    #     (
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset015_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_-_with_validation/labelsStandaloneDuplicateTs',
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset015_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_-_with_validation/predLabelsStandaloneDuplicateTs_final',
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset015_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_-_with_validation/roiMasksStandaloneDuplicateTs',
    #         False,
    #     ),
    #     (
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset015_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_-_with_validation/labelsStandaloneDuplicateVl',
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset015_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_-_with_validation/predLabelsStandaloneDuplicateVl_final',
    #         '/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset015_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_-_with_validation/roiMasksStandaloneDuplicateVl',
    #         False,
    #     ),
    # ]:
    #     write_per_tumor_analysis(gt_tumors_dir, pred_tumors_dir, rois_dir, region_based_inference)

    pass

    # def plot_corr_mat(df, title=''):
    #     corr = df.corr()
    #     plt.figure()
    #     ax = sns.heatmap(
    #         corr, annot=True,
    #         vmin=-1, vmax=1, center=0,
    #         cmap=sns.diverging_palette(20, 220, n=200),
    #         square=True
    #     )
    #     ax.set_xticklabels(
    #         ax.get_xticklabels(),
    #         rotation=45,
    #         horizontalalignment='right'
    #     )
    #     plt.title(title)
    #
    #
    # def plot_scatter(df, x, y, title=''):
    #     df.plot.scatter(x=x, y=y, c='g')
    #     plt.title(title)
    #
    #
    # def add_scan(_df):
    #     _df['scan'] = ['_-_'.join(s.split('_-_')[:-1]) for s in _df[_df.columns[0]].to_list()]
    #     return _df
    #
    # def print_summarization(gt_df, fp_pred_df, closer_than, type_column_name, dist_column_name, original_num_of_scans):
    #
    #     tp_df = gt_df[gt_df[type_column_name] == 1]
    #     fn_df = gt_df[gt_df[type_column_name] == 0]
    #     gt_closer_than_df = gt_df[gt_df[dist_column_name] <= closer_than]
    #     gt_further_than_df = gt_df[gt_df[dist_column_name] > closer_than]
    #
    #     fp_pred_closer_than_df = fp_pred_df[fp_pred_df[dist_column_name] <= closer_than]
    #     fp_pred_further_than_df = fp_pred_df[fp_pred_df[dist_column_name] > closer_than]
    #
    #     get_percentage = lambda _df: f'{_df[_df[dist_column_name] <= closer_than].shape[0]:>3}/{_df.shape[0]:<3} = {100 * _df[_df[dist_column_name] <= closer_than].shape[0] / _df.shape[0]:.2f}%'
    #
    #     print_percentage = lambda _df, type: print(f'{type:<3} lesions: {get_percentage(_df)} are closer than {closer_than} mm to the border') if _df.shape[0] > 0 else print(f'{type:<3} lesions: nan')
    #
    #     recall = lambda _df: _df[_df[type_column_name] == 1].shape[0] / _df.shape[0]
    #     precision = lambda _gt_df, _fp_pred_df: _gt_df[_gt_df[type_column_name] == 1].shape[0] / (_gt_df[_gt_df[type_column_name] == 1].shape[0] + _fp_pred_df.shape[0])
    #
    #     print_recall = lambda _df, type: print(f'{type} lesions Recall: {recall(_df):.2f}')
    #     print_precision = lambda _gt_df, _fp_pred_df, type: print(f'{type} lesions Precision: {precision(_gt_df, _fp_pred_df):.2f}')
    #
    #     def get_precision_and_recall_per_scan_averaged(_gt_df, _fp_pred_df):
    #         scans_treated = []  # in order to treat normal scans that do not have lesions in _gt_df
    #         precisions, recalls = [], []
    #         for scan_name, scan_gt_df in _gt_df.groupby('scan'):
    #             scans_treated.append(scan_name)
    #             scan_fp_pred_df = _fp_pred_df[_fp_pred_df['scan'] == scan_name]
    #             r = recall(scan_gt_df)
    #             if scan_gt_df[scan_gt_df[type_column_name] == 1].shape[0] == 0:
    #                 p = 0
    #             else:
    #                 p = precision(scan_gt_df, scan_fp_pred_df)
    #             recalls.append(r)
    #             precisions.append(p)
    #         n_scans_with_gt_lesions = len(scans_treated)
    #         # print(f'{n_scans_with_gt_lesions} scans')
    #
    #         # treating scans with only FPs lesions
    #         scans_with_only_fps = list(set(_fp_pred_df['scan'].to_list()) - set(scans_treated))
    #         n_scans_with_only_fps = len(scans_with_only_fps)
    #         precisions += [0] * n_scans_with_only_fps
    #         recalls += [1] * n_scans_with_only_fps
    #
    #         # treating normal scans perfect predictions (i.e, no TP, no FN and no FP)
    #         n_perfect_predicted_normal_scans = original_num_of_scans - n_scans_with_gt_lesions - n_scans_with_only_fps
    #         tmp = [1] * n_perfect_predicted_normal_scans
    #         precisions += tmp
    #         recalls += tmp
    #
    #         return f'{np.mean(precisions):.2f}+-{np.std(precisions):.2f}', f'{np.mean(recalls):.2f}+-{np.std(recalls):.2f}'
    #
    #     def print_precision_and_recall_per_scan_averaged(_gt_df, _fp_pred_df, type):
    #         p, r = get_precision_and_recall_per_scan_averaged(_gt_df, _fp_pred_df)
    #         print(f'{type:<25}: Precision={p}, Recall={r}')
    #
    #     # precision_per_scan_averaged, recall_per_scan_averaged = get_precision_and_recall_per_scan_averaged(gt_df, fp_pred_df)
    #
    #     print_percentage(gt_df, 'All')
    #     print_percentage(tp_df, 'TP')
    #     print_percentage(fn_df, 'FN')
    #     print_percentage(fp_pred_df, 'FP')
    #
    #     # print('')
    #     # print_recall(gt_df, 'All')
    #     # print_recall(gt_closer_than_df, f'Closer than {closer_than} mm')
    #     # print_recall(gt_further_than_df, f'Further than {closer_than} mm')
    #     #
    #     # print('')
    #     # print_precision(gt_df, fp_pred_df, 'All')
    #     # print_precision(gt_closer_than_df, fp_pred_closer_than_df, f'Closer than {closer_than} mm')
    #     # print_precision(gt_further_than_df, fp_pred_further_than_df, f'Further than {closer_than} mm')
    #
    #     print('\n########### Per Scan Averaged:')
    #     print_precision_and_recall_per_scan_averaged(gt_df, fp_pred_df, 'All lesions')
    #     print_precision_and_recall_per_scan_averaged(gt_closer_than_df, fp_pred_closer_than_df, f'Closer than {closer_than} mm lesions')
    #     print_precision_and_recall_per_scan_averaged(gt_further_than_df, fp_pred_further_than_df, f'Further than {closer_than} mm lesions')
    #
    #
    # experiment_5_pisa_ex_per_lesion_analysis = '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset005_HadassahOriginalTrainSetExperimentSimultaneousImprovedRegistration_zscore_normalization/predLabelsStandaloneDuplicatePisa_checkpoint_final/per_tumor_analysis.xlsx'
    # experiment_5_hadassah_ex_per_lesion_analysis = '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset005_HadassahOriginalTrainSetExperimentSimultaneousImprovedRegistration_zscore_normalization/predLabelsStandaloneDuplicateTs_checkpoint_final/per_tumor_analysis.xlsx'
    # experiment_8_pisa_ex_per_lesion_analysis = '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset008_HadassahOriginalTrainSetExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training/predLabelsStandaloneDuplicatePisa_checkpoint_final/per_tumor_analysis.xlsx'
    # experiment_8_hadassah_ex_per_lesion_analysis = '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset008_HadassahOriginalTrainSetExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training/predLabelsStandaloneDuplicateTs_checkpoint_final/per_tumor_analysis.xlsx'
    # experiment_9_pisa_ex_per_lesion_analysis = '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset009_HadassahOriginalTrainSetWithPisaNewTrainSetCroppedExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training/predLabelsStandaloneDuplicatePisa_checkpoint_final/per_tumor_analysis.xlsx'
    # experiment_9_hadassah_ex_per_lesion_analysis = '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset009_HadassahOriginalTrainSetWithPisaNewTrainSetCroppedExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training/predLabelsStandaloneDuplicateTs_checkpoint_final/per_tumor_analysis.xlsx'
    # experiment_10_pisa_ex_per_lesion_analysis = '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset010_HadassahOriginalTrainSetWithPisaNewTrainSetCroppedExperimentSimultaneousImprovedRegistration_zscore_normalization/predLabelsStandaloneDuplicatePisa_checkpoint_best/per_tumor_analysis.xlsx'
    # experiment_10_hadassah_ex_per_lesion_analysis = '/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset010_HadassahOriginalTrainSetWithPisaNewTrainSetCroppedExperimentSimultaneousImprovedRegistration_zscore_normalization/predLabelsStandaloneDuplicateTs_checkpoint_best/per_tumor_analysis.xlsx'
    #
    # for ex_per_lesion_analysis, experiment, dataset_name, original_num_of_scans in [
    #     (experiment_5_pisa_ex_per_lesion_analysis, 5, 'pisa', 33),
    #     (experiment_8_pisa_ex_per_lesion_analysis, 8, 'pisa', 33),
    #     (experiment_9_pisa_ex_per_lesion_analysis, 9, 'pisa', 33),
    #     (experiment_10_pisa_ex_per_lesion_analysis, 10, 'pisa', 33),
    #     (experiment_5_hadassah_ex_per_lesion_analysis, 5, 'hadassah', 33),
    #     (experiment_8_hadassah_ex_per_lesion_analysis, 8, 'hadassah', 33),
    #     (experiment_9_hadassah_ex_per_lesion_analysis, 9, 'hadassah', 33),
    #     (experiment_10_hadassah_ex_per_lesion_analysis, 10, 'hadassah', 33),
    # ]:
    #
    #     gt_df = pd.read_excel(ex_per_lesion_analysis, sheet_name='GT Lesions').iloc[:-5]
    #     fp_pred_df = pd.read_excel(ex_per_lesion_analysis, sheet_name='FP Pred Lesions').iloc[:-5]
    #
    #     gt_df = add_scan(gt_df)
    #     fp_pred_df = add_scan(fp_pred_df)
    #
    #     # gt_df['Tumor Diameter (mm)'].plot.hist(bins=100, title=dataset_name, xlabel='Tumor Diameter (mm)',
    #     #                                        xticks=list(range(1, 8, 2)) + list(range(10, 30, 4)) + list(range(30, 120, 10)))
    #     # gt_df['Distance from organ border (mm)'].plot.hist(bins=100, title=dataset_name, xlabel='Distance from organ border (mm)',
    #     #                                                    xticks=list(range(1, 8, 2)) + list(range(10, 30, 4)) + list(
    #     #                                                        range(30, 120, 10))
    #     #                                                    )
    #     #
    #     # plt.show()
    #     # continue
    #
    #     """
    #     'Unnamed: 0',
    #     'Is TP?',
    #     'Dice',
    #     'Tumor Diameter (mm)',
    #     'Distance from organ border (mm)',
    #     """
    #
    #     # def cat(r):
    #     #     v = r['Tumor Diameter (mm)']
    #     #     # if v <= 5:
    #     #     #     return '< 5'
    #     #     # elif v <= 10:
    #     #     if v <= 10:
    #     #         return '< 10'
    #     #     else:
    #     #         return '> 10'
    #     #
    #     # gt_df['cat'] = gt_df.apply(cat, axis=1)
    #
    #     # sns.lmplot(data=gt_df, x='Is TP?', y='Distance from organ border (mm)')
    #     # sns.lmplot(data=gt_df, x='Is TP?', y='Distance from organ border (mm)', hue='cat')
    #     # plt.show()
    #
    #     def filter_diameters(_df, min=0, max=np.inf):
    #         return _df[(min < _df['Tumor Diameter (mm)']) & (_df['Tumor Diameter (mm)'] <= max)]
    #
    #     possible_cats = {
    #         'lesions > 0': (0, np.inf),
    #         'lesions > 5': (5, np.inf),
    #         'lesions < 5': (0, 5),
    #         'lesions > 10': (10, np.inf),
    #         'lesions < 10': (0, 10),
    #         '5 < lesions < 10': (5, 10),
    #     }
    #
    #     # cat = 'lesions > 0'
    #     # cat = 'lesions > 5'
    #     cat = 'lesions < 5'
    #     # cat = 'lesions > 10'
    #     # cat = 'lesions < 10'
    #     # cat = '5 < lesions < 10'
    #
    #     # df_to_analyse = gt_df[gt_df['cat'] == cat]
    #
    #     cat_min, cat_max = possible_cats[cat]
    #     gt_df_to_analyse = filter_diameters(gt_df, cat_min, cat_max)
    #     fp_pred_df_to_analyse = filter_diameters(fp_pred_df, cat_min, cat_max)
    #
    #     print(f'\n\n----------------------- experiment {experiment} {dataset_name} dataset -  {cat} mm analysis ({gt_df_to_analyse.shape[0]}/{gt_df.shape[0]}= {100*gt_df_to_analyse.shape[0]/gt_df.shape[0]:.2f}% of all lesions) -----------------------\n\n')
    #     print_summarization(gt_df_to_analyse, fp_pred_df_to_analyse, closer_than=2.5, type_column_name='Is TP?',
    #                         dist_column_name='Distance from organ border (mm)',
    #                         original_num_of_scans=original_num_of_scans)
    #
    # exit(0)

    pass

    # dup_hadassa_gts = sorted(glob(f'/media/rochman/My Passport/Registered_data/*/BL_*/FU_Scan_Tumors.nii.gz'))
    # hadassa_gts, hadassa_case_ids = [], []
    # for f in dup_hadassa_gts:
    #     fu_name = f.split('/')[-2].split('_FU_')[-1]
    #     if fu_name not in hadassa_case_ids:
    #         hadassa_case_ids.append(fu_name)
    #         hadassa_gts.append(f)
    # hadassa_livers = [f.replace('/FU_Scan_Tumors.nii.gz', '/FU_Scan_Liver.nii.gz') for f in hadassa_gts]
    #
    # pisa_gts, pisa_livers = [], []
    # for folder in [
    #     'pisa_data_for_testing',
    #     'pisa_data_for_retraining',
    # ]:
    #     pisa_gts += sorted(glob(f'/home/rochman/Documents/src/data/{folder}/*/*/cropped_data/tumors.nii.gz'))
    #     pisa_livers += sorted(glob(f'/home/rochman/Documents/src/data/{folder}/*/*/cropped_data/liver_pred.nii.gz'))
    # pisa_case_ids = ['_-_'.join(f.split('/')[-4:-2]) for f in pisa_gts]
    #
    # # n_cases = 2 # todo delete
    # n_cases = None
    #
    # def extract_diameters_and_distance_from_border(gt_tumors_roi: Tuple[str, str]):
    #     tempdir = tempfile.TemporaryDirectory()
    #
    #     case, nifti_f = load_nifti_data(gt_tumors_roi[0])
    #     tmp_zeros = np.zeros_like(case)
    #     tmp_zeros_path = f'{tempdir.name}/zeros.nii.gz'
    #     nib.save(nib.Nifti1Image(tmp_zeros, nifti_f.affine, nifti_f.header), tmp_zeros_path)
    #
    #     res = analyze_per_tumor_measures((tmp_zeros_path, *gt_tumors_roi))
    #
    #     tempdir.cleanup()
    #
    #     return res
    #
    #
    # for gts, rois, case_ids, dataset_name in [
    #     (hadassa_gts, hadassa_livers, hadassa_case_ids, 'hadassa'),
    #     (pisa_gts, pisa_livers, pisa_case_ids, 'pisa'),
    # ]:
    #     gts, rois, case_ids = gts[:n_cases], rois[:n_cases], case_ids[:n_cases]
    #
    #     res = process_map(extract_diameters_and_distance_from_border, list(zip(gts, rois)), max_workers=5)
    #     # res = list(map(extract_diameters_and_distance_from_border, list(zip(gts, rois))))
    #     _, res = zip(*res)
    #
    #     res = [(f'{case_ids[j]}_-_{i+1}', *r_t) for j, r in enumerate(res) for i, r_t in enumerate(r)]
    #     columns = ['Case-ID', 'Tumor Diameter (mm)', 'Distance from organ border (mm)']
    #     res = pd.DataFrame(res, columns=columns)
    #     output_ex_fn = f'{dataset_name}.xlsx'
    #     writer = pd.ExcelWriter(output_ex_fn, engine='xlsxwriter')
    #
    #     write_to_excel(res, writer, columns, 'Case-ID')
    #
    #     writer.close()
    #
    # exit(0)

    pass

    pisa_df = pd.read_excel('/home/rochman/Documents/src/nnUNetSegmentation/measures_calculations/pisa.xlsx').iloc[:-5]
    hadassa_df = pd.read_excel('/home/rochman/Documents/src/nnUNetSegmentation/measures_calculations/hadassa.xlsx').iloc[:-5]

    # # pisa_df = pd.read_excel('/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset009_HadassahOriginalTrainSetWithPisaNewTrainSetCroppedExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training/predLabelsStandaloneDuplicatePisa_checkpoint_final/per_tumor_analysis.xlsx').iloc[:-5]
    # # hadassa_df = pd.read_excel('/home/rochman/Documents/src/data/nnUNet/nnUNet_inference/Dataset009_HadassahOriginalTrainSetWithPisaNewTrainSetCroppedExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training/predLabelsStandaloneDuplicateTs_checkpoint_final/per_tumor_analysis.xlsx'
    # #     ).iloc[:-5]

    def plot_hist(df, dataset_name, column, bins=100, min=None, max=None, percentage=False, cumulative=False):

        # if percentage:
        #     plt.figure()

        df_to_plot = df
        if min is not None:
            df_to_plot = df_to_plot[df_to_plot[column] >= min]
        if max is not None:
            df_to_plot = df_to_plot[df_to_plot[column] < max]

        if min is None and max is None:
            title = 'All lesions'
        elif min is None and max is not None:
            title = f'Lesions < {max}'
        elif min is not None and max is None:
            title = f'Lesions > {min}'
        else:
            title = f'{min} < Lesions < {max}'

        if cumulative:
            title = f'{title} (cumulative)'

        plt.hist(df_to_plot[column], bins=bins, label=dataset_name,
                 weights=np.ones(len(df_to_plot)) / len(df_to_plot) if percentage else None,
                 cumulative=cumulative, alpha=0.5
                 )
        plt.xlabel(column)
        plt.title(title)

        if percentage:
            plt.gca().yaxis.set_major_formatter(PercentFormatter(1))

    percentage = False
    cumulative = True
    bins = 1000
    min = None
    max = None
    column = 'Tumor Diameter (mm)'
    # plot_hist(hadassa_df, 'hadassa', 'Tumor Diameter (mm)', bins, min, max, percentage=percentage, cumulative=cumulative)
    # plt.legend()
    # plot_hist(pisa_df, 'pisa', 'Tumor Diameter (mm)', bins, min, max, percentage=percentage, cumulative=cumulative)
    # plt.legend()
    # plt.show()

    # plot_hist(hadassa_df, 'hadassa', 'Distance from organ border (mm)', bins, min, max, percentage=percentage,
    #           cumulative=cumulative)
    # plt.legend()
    # plot_hist(pisa_df, 'pisa', 'Distance from organ border (mm)', bins, min, max, percentage=percentage, cumulative=cumulative)
    # plt.legend()
    # plt.show()

    pass

    # fold_1_test_cases = [f.split('/')[-1].replace('.nii.gz', '') for f in glob(
    #     '/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset011_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_fold_1/labelsStandaloneDuplicateTs/*.nii.gz')]
    # fold_2_test_cases = [f.split('/')[-1].replace('.nii.gz', '') for f in glob(
    #     '/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset012_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_fold_2/labelsStandaloneDuplicateTs/*.nii.gz')]
    # fold_3_test_cases = [f.split('/')[-1].replace('.nii.gz', '') for f in glob(
    #     '/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset013_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_fold_3/labelsStandaloneDuplicateTs/*.nii.gz')]
    #
    # df = pd.concat([pisa_df, hadassa_df])
    #
    # df['case_name'] = df.apply(lambda r: '_-_'.join(r[df.columns[0]].split('_-_')[:-1]), axis=1)
    #
    # df_fold_1_test = df[df['case_name'].isin(fold_1_test_cases)]
    # df_fold_2_test = df[df['case_name'].isin(fold_2_test_cases)]
    # df_fold_3_test = df[df['case_name'].isin(fold_3_test_cases)]
    #
    # print(f'1st experiment # of lesions: train {len(df_fold_2_test) + len(df_fold_3_test)}, test {len(df_fold_1_test)}')
    # print(f'2nd experiment # of lesions: train {len(df_fold_1_test) + len(df_fold_3_test)}, test {len(df_fold_2_test)}')
    # print(f'3rd experiment # of lesions: train {len(df_fold_2_test) + len(df_fold_1_test)}, test {len(df_fold_3_test)}')
    #
    # for fold, (test_df, train_dfs) in enumerate([
    #     (
    #         df_fold_1_test, (df_fold_2_test, df_fold_3_test)
    #     ),
    #     (
    #         df_fold_2_test, (df_fold_1_test, df_fold_3_test)
    #     ),
    #     (
    #         df_fold_3_test, (df_fold_2_test, df_fold_1_test)
    #     ),
    # ]):
    #     train_df = pd.concat(train_dfs)
    #
    #     plt.figure()
    #     plot_hist(train_df, f'train_fold_{fold + 1}', 'Tumor Diameter (mm)', bins, min, max, percentage=percentage, cumulative=cumulative)
    #     plt.legend()
    #     plot_hist(test_df, f'test_fold_{fold + 1}', 'Tumor Diameter (mm)', bins, min, max, percentage=percentage, cumulative=cumulative)
    #     plt.legend()
    #
    # plt.show()

    pass

    test_cases = [f.split('/')[-1].replace('.nii.gz', '') for f in glob('/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/labelsStandaloneDuplicateTs/*.nii.gz')]
    val_cases = [f.split('/')[-1].replace('.nii.gz', '') for f in glob('/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/labelsStandaloneDuplicateVl/*.nii.gz')]
    train_and_val_cases = [f.split('/')[-1].replace('.nii.gz', '').split('_FU_')[-1] for f in glob('/home/rochman/Documents/src/data/nnUNet/nnUNet_raw/Dataset014_AllDataExperimentSimultaneousImprovedRegistration_zscore_normalization_region_based_training_-_with_validation/labelsTr/*.nii.gz')]
    train_cases = list(set(train_and_val_cases) - set(val_cases))

    df = pd.concat([pisa_df, hadassa_df])

    df['case_name'] = df.apply(lambda r: '_-_'.join(r[df.columns[0]].split('_-_')[:-1]), axis=1)
    df['patient'] = df.apply(lambda r: r[df.columns[0]].split('_-_')[0] if (r[df.columns[0]].count('_-_') == 2) else '_'.join(c for c in r[df.columns[0]].split('_-_')[0].split('_') if not c.isdigit()), axis=1)

    df_test = df[df['case_name'].isin(test_cases)]
    df_val = df[df['case_name'].isin(val_cases)]
    df_train = df[df['case_name'].isin(train_cases)]

    print(f'# of lesions: train {len(df_train)}, val {len(df_val)}, test {len(df_test)}')

    for (set_df, set_name) in [
        (
            df_train, 'train'
        ),
        (
            df_val, 'val'
        ),
        (
            df_test, 'test'
        ),
    ]:

        plot_hist(set_df, set_name, column, bins, min, max, percentage=percentage, cumulative=cumulative)

        print(f'{set_name} Mean: {set_df[column].mean():.2f}, STD: {set_df[column].std():.2f}')

    plt.legend()
    plt.show()

    pass

    # df = pd.concat([pisa_df, hadassa_df])
    #
    # from sklearn.model_selection import StratifiedGroupKFold
    #
    # df['patient'] = df.apply(lambda r: r[df.columns[0]].split('_-_')[0] if (r[df.columns[0]].count('_-_') == 2) else '_'.join(c for c in r[df.columns[0]].split('_-_')[0].split('_') if not c.isdigit()), axis=1)
    # df['diameter_bins'] = np.digitize(df['Tumor Diameter (mm)'], np.arange(4, 25))
    # df['distance_bins'] = np.digitize(df['Distance from organ border (mm)'], np.arange(1, 15, 2))
    #
    # seed = 9
    #
    # train_and_val_index, test_index = next(StratifiedGroupKFold(n_splits=5, shuffle=True, random_state=seed).split(
    #     X=df[df.columns[0]], y=df['diameter_bins'], groups=df['patient']))
    #
    # df_test = df.iloc[test_index]
    # train_and_val_df = df.iloc[train_and_val_index]
    #
    # train_index, val_index = next(StratifiedGroupKFold(n_splits=8, shuffle=True, random_state=seed).split(
    #     X=train_and_val_df[df.columns[0]], y=train_and_val_df['diameter_bins'], groups=train_and_val_df['patient']))
    #
    # df_train = train_and_val_df.iloc[train_index]
    # df_val = train_and_val_df.iloc[val_index]
    #
    # # print(len(df_train)/len(df))
    # # print(len(df_test)/len(df))
    # # print(len(df_val)/len(df))
    # # exit(0)
    #
    # print(f'# of lesions: train {len(df_train)}, val {len(df_val)}, test {len(df_test)}')
    #
    # for (set_df, set_name) in [
    #     (
    #         df_train, 'train'
    #     ),
    #     (
    #         df_val, 'val'
    #     ),
    #     (
    #         df_test, 'test'
    #     ),
    # ]:
    #
    #     plot_hist(set_df, set_name, column, bins, min, max, percentage=percentage, cumulative=cumulative)
    #
    #     print(f'{set_name} Mean: {set_df[column].mean():.2f}, STD: {set_df[column].std():.2f}')
    #
    # plt.legend()
    # plt.show()
